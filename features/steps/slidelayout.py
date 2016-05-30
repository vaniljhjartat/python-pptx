# encoding: utf-8

"""
Step implementations for slide layout-related features
"""

from __future__ import absolute_import

from behave import given, then

from pptx import Presentation
from pptx.shapes.base import BaseShape
from pptx.shapes.picture import Picture
from pptx.shapes.placeholder import LayoutPlaceholder
from pptx.shapes.shapetree import LayoutPlaceholders, LayoutShapes

from helpers import test_pptx


SHAPE_COUNT = 3


# given ===================================================

@given('a layout placeholder collection')
def given_layout_placeholder_collection(context):
    prs = Presentation(test_pptx('lyt-shapes'))
    context.layout_placeholders = prs.slide_layouts[0].placeholders


@given('a layout shape collection')
def given_layout_shape_collection(context):
    prs = Presentation(test_pptx('lyt-shapes'))
    context.layout_shapes = prs.slide_layouts[0].shapes


@given('a slide layout')
def given_a_slide_layout(context):
    prs = Presentation(test_pptx('sld-slide-access'))
    context.slide_layout = prs.slide_layouts[0]


@given('a slide layout having three shapes')
def given_slide_layout_having_three_shapes(context):
    prs = Presentation(test_pptx('lyt-shapes'))
    context.slide_layout = prs.slide_layouts[0]


@given('a slide layout having two placeholders')
def given_layout_having_two_placeholders(context):
    prs = Presentation(test_pptx('lyt-shapes'))
    context.slide_layout = prs.slide_layouts[0]


# then ====================================================

@then('each shape is of the appropriate type')
def then_each_shape_is_of_appropriate_type(context):
    layout_shapes = context.layout_shapes
    expected_types = [LayoutPlaceholder, LayoutPlaceholder, Picture]
    for idx, layout_shape in enumerate(layout_shapes):
        assert type(layout_shape) == expected_types[idx], (
            "got \'%s\'" % type(layout_shape).__name__
        )


@then('I can access a layout placeholder by idx value')
def then_can_access_layout_placeholder_by_idx_value(context):
    layout_placeholders = context.layout_placeholders
    title_placeholder = layout_placeholders.get(idx=0)
    body_placeholder = layout_placeholders.get(idx=10)
    assert title_placeholder._element is layout_placeholders[0]._element
    assert body_placeholder._element is layout_placeholders[1]._element


@then('I can access a layout placeholder by index')
def then_can_access_layout_placeholder_by_index(context):
    layout_placeholders = context.layout_placeholders
    for idx in range(2):
        layout_placeholder = layout_placeholders[idx]
        assert isinstance(layout_placeholder, LayoutPlaceholder)


@then('I can access a layout shape by index')
def then_can_access_layout_shape_by_index(context):
    layout_shapes = context.layout_shapes
    for idx in range(SHAPE_COUNT):
        layout_shape = layout_shapes[idx]
        assert isinstance(layout_shape, BaseShape)


@then('I can access the placeholder collection of the slide layout')
def then_can_access_placeholder_collection_of_slide_layout(context):
    slide_layout = context.slide_layout
    layout_placeholders = slide_layout.placeholders
    msg = 'SlideLayout.placeholders not instance of LayoutPlaceholders'
    assert isinstance(layout_placeholders, LayoutPlaceholders), msg


@then('I can access the shape collection of the slide layout')
def then_can_access_shape_collection_of_slide_layout(context):
    slide_layout = context.slide_layout
    layout_shapes = slide_layout.shapes
    msg = 'SlideLayout.shapes not instance of LayoutShapes'
    assert isinstance(layout_shapes, LayoutShapes), msg


@then('I can iterate over the layout placeholders')
def then_can_iterate_over_the_layout_placeholders(context):
    layout_placeholders = context.layout_placeholders
    actual_count = 0
    for layout_placeholder in layout_placeholders:
        actual_count += 1
        assert isinstance(layout_placeholder, LayoutPlaceholder)
    assert actual_count == 2


@then('I can iterate over the layout shapes')
def then_can_iterate_over_the_layout_shapes(context):
    layout_shapes = context.layout_shapes
    actual_count = 0
    for layout_shape in layout_shapes:
        actual_count += 1
        assert isinstance(layout_shape, BaseShape)
    assert actual_count == SHAPE_COUNT


@then('slide_layout.slide_master is a SlideMaster object')
def then_slide_layout_slide_master_is_a_SlideMaster_object(context):
    slide_layout = context.slide_layout
    assert type(slide_layout.slide_master).__name__ == 'SlideMaster'


@then('slide_layout.placeholders is a LayoutPlaceholders object')
def then_slide_layout_placeholders_is_a_LayoutPlaceholders_object(context):
    slide_layout = context.slide_layout
    assert type(slide_layout.placeholders).__name__ == 'LayoutPlaceholders'


@then('slide_layout.shapes is a LayoutShapes object')
def then_slide_layout_shapes_is_a_LayoutShapes_object(context):
    slide_layout = context.slide_layout
    assert type(slide_layout.shapes).__name__ == 'LayoutShapes'


@then('the length of the layout placeholder collection is 2')
def then_len_of_placeholder_collection_is_2(context):
    slide_layout = context.slide_layout
    layout_placeholders = slide_layout.placeholders
    assert len(layout_placeholders) == 2, (
        'expected len(layout_placeholders) of 2, got %s' %
        len(layout_placeholders)
    )


@then('the length of the layout shape collection counts all its shapes')
def then_len_of_layout_shape_collection_counts_all_its_shapes(context):
    slide_layout = context.slide_layout
    layout_shapes = slide_layout.shapes
    assert len(layout_shapes) == SHAPE_COUNT, (
        'expected len(layout_shapes) of %d, got %s' %
        (SHAPE_COUNT, len(layout_shapes))
    )
